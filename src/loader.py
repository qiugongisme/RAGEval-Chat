import os
import re
from typing import List, Optional

import fitz  # PyMuPDF
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import config


def get_pdf_files(directory: str) -> List[str]:
    """获取指定目录下的所有PDF文件路径
    - param directory: 目录路径
    """
    pdf_files = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith('.pdf'):
                pdf_files.append(os.path.join(root, file))
    return pdf_files


def save_text_to_file(text: str, output_path: str) -> None:
    """将提取的文本保存到指定文件"""
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as file:
        file.write(text)


def rm_useless_content(text):
    """去除页眉页脚等不必要的内容"""
    fixed_strings_to_remove = [
        "国家金融监督管理总局规章",
        "国家金融监督管理总局发布",
        "中国银行保险监督管理委员会规章",
        "中国银行保险监督管理委员会发布"
    ]
    for item in fixed_strings_to_remove:
        text = text.replace(item, "")
    # 去除 "- 数字 -" 格式的页码
    text = re.sub(r"-\s*\d+\s*-", "", text)
    # 去除网址
    text = re.sub(r'(https?://[^\s]+)', '', text)
    # 去除 docId=xxx&itemId=yyy
    text = re.sub(r'docId=\d+(&itemId=\d+)?', '', text)
    # 去除多余空行
    text = re.sub(r'\n\s*\n', '\n', text).strip()
    return text


def load_pdf2document(directory: str) -> List[Document]:
    """从指定目录加载 PDF 文件并提取文本内容"""
    pdf_file_list = get_pdf_files(directory)
    texts, metadatas = [], []
    for pdf_file in pdf_file_list:
        document = fitz.open(pdf_file)
        text = ""
        for page_num in range(len(document)):
            page = document.load_page(page_num)
            text += page.get_text()

        text = rm_useless_content(text)

        output_path = os.path.join(config.FILE_OUTPUT_PATH, os.path.basename(pdf_file).replace('.pdf', '.txt'))
        save_text_to_file(text, output_path)

        texts.append(text)
        metadatas.append({"source": "《" + os.path.basename(pdf_file).replace('.pdf', '') + "》"})

    documents = RecursiveCharacterTextSplitter().create_documents(texts, metadatas=metadatas)
    return documents


def _extract_text_from_pdf(filepath: str) -> str:
    """提取 PDF 文件文本"""
    doc = fitz.open(filepath)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text


def _extract_text_from_docx(filepath: str) -> str:
    """提取 DOCX 文件文本"""
    from docx import Document
    doc = Document(filepath)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_text_from_doc(filepath: str) -> str:
    """提取 DOC 文件文本（旧格式，尝试用 python-docx 读取）"""
    try:
        return _extract_text_from_docx(filepath)
    except Exception:
        # python-docx 不支持旧 .doc 格式，fallback 到原始文本提取
        try:
            with open(filepath, 'rb') as f:
                raw = f.read()
            # 尝试提取可读文本
            text = raw.decode('utf-8', errors='ignore')
            # 过滤掉二进制垃圾字符
            text = re.sub(r'[^\u4e00-\u9fff\u0020-\u007e\n]', '', text)
            return text.strip()
        except Exception:
            return ""


def _extract_text_from_pptx(filepath: str) -> str:
    """提取 PPTX 文件文本"""
    from pptx import Presentation
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


def _extract_text_from_ppt(filepath: str) -> str:
    """提取 PPT 文件文本（旧格式，尝试用 python-pptx 读取）"""
    try:
        return _extract_text_from_pptx(filepath)
    except Exception:
        try:
            with open(filepath, 'rb') as f:
                raw = f.read()
            text = raw.decode('utf-8', errors='ignore')
            text = re.sub(r'[^\u4e00-\u9fff\u0020-\u007e\n]', '', text)
            return text.strip()
        except Exception:
            return ""


def _extract_text_from_txt(filepath: str) -> str:
    """提取 TXT 文件文本"""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _extract_text_from_md(filepath: str) -> str:
    """提取 MD 文件文本"""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


SUPPORTED_EXTENSIONS = {
    '.pdf': _extract_text_from_pdf,
    '.docx': _extract_text_from_docx,
    '.doc': _extract_text_from_doc,
    '.pptx': _extract_text_from_pptx,
    '.ppt': _extract_text_from_ppt,
    '.txt': _extract_text_from_txt,
    '.md': _extract_text_from_md,
}


def load_single_document(filepath: str) -> Optional[Document]:
    """加载单个文档文件，返回 Document 对象"""
    ext = os.path.splitext(filepath)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return None

    extractor = SUPPORTED_EXTENSIONS[ext]
    text = extractor(filepath)
    if not text.strip():
        return None

    # PDF 特殊处理：去除无用内容和保存 txt
    if ext == '.pdf':
        text = rm_useless_content(text)
        output_path = os.path.join(config.FILE_OUTPUT_PATH, os.path.basename(filepath).replace('.pdf', '.txt'))
        save_text_to_file(text, output_path)

    source_name = os.path.splitext(os.path.basename(filepath))[0]

    return Document(
        page_content=text,
        metadata={"source": "《" + source_name + "》", "filepath": filepath, "type": ext[1:]}
    )


def load_documents_from_paths(file_paths: List[str]) -> List[Document]:
    """从文件路径列表加载所有支持的文档"""
    documents = []
    for fp in file_paths:
        if not os.path.isfile(fp):
            continue
        doc = load_single_document(fp)
        if doc:
            documents.append(doc)
    return documents


def split_by_pattern(content: str, pattern: str = r"第\S*条") -> List[str]:
    """根据正则表达式切分内容"""
    matches = list(re.finditer(rf"^{pattern}", content, re.MULTILINE))
    if not matches:
        return [content.strip()]

    result = []
    for i, match in enumerate(matches):
        start = match.start()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(content)
        part = content[start:end].strip()
        if part:
            result.append(part)
    return result


class CustomDocument:
    def __init__(self, content: str, metadata: dict) -> None:
        self.content = content
        self.metadata = metadata


def load_and_split(directory: str) -> list[CustomDocument]:
    """从指定文件目录加载 PDF 文件并提取、切分文本内容"""
    result = []
    pdf_file_list = get_pdf_files(directory)
    for pdf_file in pdf_file_list:
        document = fitz.open(pdf_file)
        text_content = ""
        for page_num in range(len(document)):
            page = document.load_page(page_num)
            text_content += page.get_text()

        text_content = rm_useless_content(text_content)

        output_path = os.path.join(config.FILE_OUTPUT_PATH, os.path.basename(pdf_file).replace('.pdf', '.txt'))
        save_text_to_file(text_content, output_path)

        split_list = split_by_pattern(text_content)
        metadata = {"source": "《" + os.path.basename(pdf_file).replace('.pdf', '') + "》"}
        for split_content in split_list:
            result.append(CustomDocument(split_content, metadata))

    return result
